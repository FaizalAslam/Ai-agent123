from flask import Flask, render_template, request, jsonify
import threading
import json
import time
import logging
import webbrowser
import traceback
import os
import re
from pathlib import Path

# ---- Core modules ---------------------------------------------------------
from modules import system_core, ui, config

# ---- Office Agent (Project 2) --------------------------------------------
from utils.command_buffer import CommandBuffer
from utils import command_map
from utils.office_actions import OfficeActionError, normalize_actions, validate_actions
from executor.excel_executor import ExcelExecutor
from executor.word_executor import WordExecutor
from executor.ppt_executor import PowerPointExecutor
from parser.command_parser import parse_command
from ai.openai_handler import OpenAIHandler
from listener.keyboard_listener import KeyboardListener
from listener.clipboard_listener import ClipboardListener
try:
    from listener.voice_listener import VoiceListener
    VOICE_MODULE_AVAILABLE = True
except Exception:
    VoiceListener = None
    VOICE_MODULE_AVAILABLE = False

# ---- Optional modules (graceful fallback) --------------------------------
try:
    from modules import ocr_utils
    OCR_AVAILABLE = True
except Exception as e:
    print(f"OCR unavailable: {e}")
    OCR_AVAILABLE = False

try:
    from modules import pdf_utils
    PDF_AVAILABLE = True
except Exception as e:
    print(f"PDF unavailable: {e}")
    PDF_AVAILABLE = False

try:
    from modules import pdf_reader
    READER_AVAILABLE = True
except Exception as e:
    print(f"PDF Reader unavailable: {e}")
    READER_AVAILABLE = False

try:
    from modules import gui_automation
    GUI_AVAILABLE = True
except Exception as e:
    print(f"GUI unavailable: {e}")
    GUI_AVAILABLE = False

try:
    from modules import pdf_editor
    PDF_EDITOR_AVAILABLE = True
except Exception as e:
    print(f"PDF Editor unavailable: {e}")
    PDF_EDITOR_AVAILABLE = False

try:
    import keyboard
    KEYBOARD_AVAILABLE = True
except ImportError:
    print("keyboard not found — pip install keyboard")
    KEYBOARD_AVAILABLE = False

# ---- Logging --------------------------------------------------------------
logging.basicConfig(
    filename="agent.log",
    level=logging.INFO,
    format="%(asctime)s - %(message)s",
    datefmt="%H:%M:%S",
    filemode="w"
)
logging.getLogger("werkzeug").setLevel(logging.ERROR)

# ---- Flask app ------------------------------------------------------------
app = Flask(__name__)

# ---- Shared state ---------------------------------------------------------
last_ocr = {"text": "", "pending": False}

# ---- Office Agent setup ---------------------------------------------------
OFFICE_APPS = {"excel", "word", "powerpoint", "ppt"}
BASE_DIR = Path(__file__).resolve().parent
OFFICE_OUTPUT_DIR = BASE_DIR / "outputs"
OFFICE_EXTENSIONS = {
    "excel": "xlsx",
    "word": "docx",
    "powerpoint": "pptx",
    "ppt": "pptx",
}
OFFICE_OUTPUT_PREFIXES = {
    "excel": "excel_output",
    "word": "word_output",
    "powerpoint": "powerpoint_output",
    "ppt": "powerpoint_output",
}
OFFICE_OUTPUTS = {
    "excel": str(OFFICE_OUTPUT_DIR / "output.xlsx"),
    "word": str(OFFICE_OUTPUT_DIR / "output.docx"),
    "powerpoint": str(OFFICE_OUTPUT_DIR / "output.pptx"),
    "ppt": str(OFFICE_OUTPUT_DIR / "output.pptx"),
}
OFFICE_DEPENDENCIES = {
    "excel": ("openpyxl", "openpyxl"),
    "word": ("docx", "python-docx"),
    "powerpoint": ("pptx", "python-pptx"),
    "ppt": ("pptx", "python-pptx"),
}
OFFICE_TARGET_KEYWORDS = {
    "excel": (
        "excel", "spreadsheet", "workbook", "worksheet", "sheet", "xlsx"
    ),
    "word": (
        "word", "document", "docx"
    ),
    "powerpoint": (
        "powerpoint", "power point", "ppt", "pptx", "presentation",
        "slide deck", "slides", "slide"
    ),
}
OFFICE_ACTION_KEYWORDS = (
    "create", "make", "generate", "build", "new", "add", "insert", "edit",
    "update", "modify", "write", "format", "table", "chart", "row", "column",
    "cell", "slide", "paragraph", "heading", "workbook", "worksheet",
    "spreadsheet", "document", "presentation"
)
APP_LAUNCH_PREFIXES = ("open ", "launch ", "start ", "run ", "boot ")
_cmd_buf = CommandBuffer()
_clipboard_listener = ClipboardListener(_cmd_buf)
_keyboard_listener = KeyboardListener(_handle_global_command := None, _cmd_buf)
_voice_listener = VoiceListener(_handle_global_command) if VOICE_MODULE_AVAILABLE else None
voice_state = {"enabled": False}
_openai_handler = OpenAIHandler()


# ---- Office Agent helpers -------------------------------------------------
def _safe_speak(text):
    try:
        ui.speak(text)
    except Exception:
        pass


def _extract_office_agent_command(raw_text):
    text = (raw_text or "").strip()
    match = re.match(r"^agent\s*:\s*(excel|word|powerpoint|ppt)\s*:\s*(.+)$", text, re.IGNORECASE)
    if not match:
        return None, None
    app_name = match.group(1).lower().strip()
    command_text = match.group(2).strip()
    return app_name, command_text


def _canonical_office_app(app_name):
    app_name = (app_name or "").lower().strip()
    return "powerpoint" if app_name == "ppt" else app_name


def _contains_term(text, term):
    term = (term or "").lower().strip()
    if not term:
        return False
    if " " in term:
        return term in text
    return bool(re.search(rf"\b{re.escape(term)}\b", text))


def _detect_office_intent(raw_text):
    original = (raw_text or "").strip()
    parsed_app, parsed_command = _extract_office_agent_command(original)
    if parsed_app and parsed_command:
        app = _canonical_office_app(parsed_app)
        return {
            "is_office": True,
            "app_type": app,
            "command": parsed_command,
            "action_type": _detect_action_type(parsed_command, []),
            "reason": "agent-prefixed office command",
        }

    text = re.sub(r"\s+", " ", original.lower()).strip()
    if not text:
        return {"is_office": False, "reason": "empty command"}

    app_type = None
    matched_term = ""
    for candidate, terms in OFFICE_TARGET_KEYWORDS.items():
        for term in terms:
            if _contains_term(text, term):
                app_type = candidate
                matched_term = term
                break
        if app_type:
            break

    if not app_type:
        return {"is_office": False, "reason": "no office target keyword"}

    has_action_term = any(_contains_term(text, term) for term in OFFICE_ACTION_KEYWORDS)
    has_open_doc_term = (
        text.startswith(("open ", "load ", "import "))
        and any(_contains_term(text, term) for term in (
            "file", "workbook", "worksheet", "spreadsheet", "document",
            "docx", "xlsx", "pptx", "presentation", "slide deck"
        ))
    )

    if has_action_term or has_open_doc_term:
        return {
            "is_office": True,
            "app_type": app_type,
            "command": original,
            "action_type": _detect_action_type(original, []),
            "reason": f"office target '{matched_term}' with document action term",
        }

    return {
        "is_office": False,
        "app_type": app_type,
        "reason": f"office app launch only for '{matched_term}'",
    }


def _is_app_launch_command(command_text):
    return (command_text or "").lower().strip().startswith(APP_LAUNCH_PREFIXES)


def _is_known_office_app(app_name):
    return _canonical_office_app(app_name) in {"excel", "word", "powerpoint"}


def _resolve_actions(app_name, command_text):
    def _estimate_subcommands(text):
        protected = re.sub(
            r"(\d+\s*(?:col|cols|column|columns)\s+)and(\s*\d+\s*(?:row|rows))",
            r"\1__AND__\2",
            (text or "").lower().strip()
        )
        protected = re.sub(
            r"(\d+\s*(?:row|rows)\s+)and(\s*\d+\s*(?:col|cols|column|columns))",
            r"\1__AND__\2",
            protected
        )
        parts = re.split(r"\s+(?:and|then|also|after that|next)\s+", protected)
        return max(1, len([p for p in parts if p.strip()]))

    def _actions_cover_command_intents(app, text, actions):
        low = (text or "").lower()
        names = {
            str(a.get("action", "")).strip().lower()
            for a in (actions or [])
            if isinstance(a, dict)
        }
        if not names:
            return False

        checks = []
        if app == "excel":
            if "background color" in low:
                checks.append("set_bg_color" in names)
            if "font color" in low:
                checks.append("set_font_color" in names)
            if "font size" in low:
                checks.append("set_font_size" in names)
            if "number format" in low:
                checks.append("set_number_format" in names)
            if "formula" in low or "sum" in low:
                checks.append("write_formula" in names)
            if "rename" in low and "sheet" in low:
                checks.append("rename_sheet" in names)
            if "insert row" in low:
                checks.append("insert_row" in names)
            if "[[" in low and "write" in low and "values" in low:
                checks.append("write_range" in names)

            # If command explicitly targets multiple cells for background color,
            # ensure cached actions include multiple bg actions.
            if "background color" in low and "cells" in low and re.search(r"\b[A-Z]{1,3}\d{1,7}\b.*\band\b.*\b[A-Z]{1,3}\d{1,7}\b", text, re.IGNORECASE):
                bg_count = sum(1 for a in (actions or []) if isinstance(a, dict) and str(a.get("action", "")).lower() == "set_bg_color")
                checks.append(bg_count >= 2)

        return all(checks) if checks else True

    cache_key, cached_actions, cache_score = command_map.get_cached_actions(app_name, command_text)
    # Use cache only for exact matches; fuzzy cache reuse can apply stale actions
    # to similar-but-different commands.
    if cached_actions and cache_score == 100:
        cached_count = len([a for a in cached_actions if isinstance(a, dict) and a.get("action")])
        clause_count = _estimate_subcommands(command_text)
        if cached_count >= clause_count and _actions_cover_command_intents(app_name, command_text, cached_actions):
            logging.info(f"Office cache hit [{app_name}] score={cache_score}: {command_text}")
            return cache_key or command_text, cached_actions, "command-cache", None
        logging.info(
            f"Ignoring stale cache for [{app_name}] command (cached={cached_count}, clauses={clause_count}): {command_text}"
        )

    actions = parse_command(app_name, command_text)
    if actions:
        try:
            actions = normalize_actions(actions)
        except OfficeActionError as exc:
            return command_text, [], "json-parser", exc
        # If parser returns fewer actions than apparent command clauses,
        # try API and prefer the richer valid result.
        clause_count = _estimate_subcommands(command_text)
        if clause_count > len(actions):
            ai_actions = _openai_handler.interpret(app_name, command_text)
            try:
                normalized_ai = normalize_actions(ai_actions) if ai_actions else []
            except OfficeActionError as exc:
                return command_text, [], "openai-fallback", exc
            if normalized_ai:
                if len(normalized_ai) >= len(actions):
                    command_map.save_actions(app_name, command_text, normalized_ai)
                    return command_text, normalized_ai, "openai-fallback", None
        command_map.save_actions(app_name, command_text, actions)
        return command_text, actions, "json-parser", None

    ai_actions = _openai_handler.interpret(app_name, command_text)
    try:
        normalized = normalize_actions(ai_actions) if ai_actions else []
    except OfficeActionError as exc:
        return command_text, [], "openai-fallback", exc
    if normalized:
        command_map.save_actions(app_name, command_text, normalized)
        return command_text, normalized, "openai-fallback", None

    if getattr(_openai_handler, "last_error_code", "") in {"INVALID_OPENAI_JSON", "INVALID_OFFICE_ACTION"}:
        return command_text, [], "openai-fallback", OfficeActionError(
            getattr(_openai_handler, "last_error_code", "INVALID_OPENAI_JSON"),
            "OpenAI returned an invalid Office action plan.",
            getattr(_openai_handler, "last_error", ""),
        )

    fallback = _default_create_action(app_name, command_text)
    if fallback:
        return command_text, [fallback], "office-intent-fallback", None

    return command_text, [], "no-match", None


def _extract_named_file_path(command_text, app_name):
    text = (command_text or "").strip()
    ext = {
        "excel": "xlsx",
        "word": "docx",
        "powerpoint": "pptx",
        "ppt": "pptx",
    }.get(app_name)
    if not text or not ext:
        return ""

    def _command_path(raw_name):
        path = Path((raw_name or "").strip())
        if path.is_absolute():
            return str(path.resolve())
        return str((OFFICE_OUTPUT_DIR / path).resolve())

    def _sanitize_base(name):
        cleaned = re.sub(r'[<>:"/\\|?*]+', "", (name or "").strip())
        cleaned = re.split(r"\s+(?:and|then|with|in|on)\b", cleaned, maxsplit=1, flags=re.IGNORECASE)[0]
        cleaned = re.sub(r"\s+", " ", cleaned).strip(" .")
        return cleaned

    quoted = re.search(r'["\']([^"\']+\.' + re.escape(ext) + r')["\']', text, re.IGNORECASE)
    if quoted:
        return _command_path(quoted.group(1).strip())

    plain = re.search(r'\b([A-Za-z0-9_\- .]+\.' + re.escape(ext) + r')\b', text, re.IGNORECASE)
    if plain:
        return _command_path(plain.group(1).strip())

    # Support "named demo" or "called demo" without extension.
    named = re.search(r'\b(?:named|called|name)\s*[:=]?\s*["\']?([A-Za-z0-9_\- ]{1,100})["\']?\b', text, re.IGNORECASE)
    if named:
        base = _sanitize_base(named.group(1))
        if base:
            return _command_path(f"{base}.{ext}")

    return ""


def _next_available_path(path):
    base, ext = os.path.splitext(os.path.abspath(path))
    candidate = f"{base}{ext}"
    idx = 1
    while os.path.exists(candidate):
        candidate = f"{base}_{idx}{ext}"
        idx += 1
    return candidate


def _generate_new_output_path(app_name):
    OFFICE_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ext = OFFICE_EXTENSIONS.get(app_name, "xlsx")
    prefix = OFFICE_OUTPUT_PREFIXES.get(app_name, f"{app_name}_output")
    stamp = time.strftime("%Y%m%d_%H%M%S")
    millis = int((time.time() * 1000) % 1000)
    return str((OFFICE_OUTPUT_DIR / f"{prefix}_{stamp}_{millis:03d}.{ext}").resolve())


def _resolve_path_value(value, app_name, for_output=False):
    raw = str(value or "").strip().strip('"').strip("'")
    if not raw:
        return ""

    ext = OFFICE_EXTENSIONS.get(app_name, "")
    expanded = os.path.expandvars(os.path.expanduser(raw))
    path = Path(expanded)
    if ext and not path.suffix and for_output:
        path = path.with_suffix(f".{ext}")
    if not path.is_absolute():
        base = OFFICE_OUTPUT_DIR if for_output and len(path.parts) == 1 else BASE_DIR
        path = base / path
    try:
        return str(path.resolve())
    except OSError:
        return str(path.absolute())


def _first_action_path(actions, action_names, path_keys=("path", "file_path", "filename", "output_path")):
    for action in actions or []:
        if not isinstance(action, dict):
            continue
        if str(action.get("action", "")).strip().lower() not in action_names:
            continue
        for key in path_keys:
            value = action.get(key)
            if str(value or "").strip():
                return str(value).strip(), action
    return "", None


def _save_as_action_names(app_name):
    return {
        "excel": {"save_workbook_as"},
        "word": {"save_document_as"},
        "powerpoint": {"save_presentation_as"},
        "ppt": {"save_presentation_as"},
    }.get(app_name, set())


def _open_action_names(app_name):
    return {
        "excel": {"open_workbook"},
        "word": {"open_document"},
        "powerpoint": {"open_presentation"},
        "ppt": {"open_presentation"},
    }.get(app_name, set())


def resolve_office_file_path(request_payload, actions, app_type, mode=None):
    app_name = _canonical_office_app(app_type)
    actions = actions or []
    command_text = (
        (request_payload or {}).get("raw")
        or (request_payload or {}).get("command")
        or ""
    )
    action_type = mode or _detect_action_type(command_text, actions)
    OFFICE_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    explicit = (
        (request_payload or {}).get("file_path")
        or (request_payload or {}).get("file")
        or ""
    )
    explicit_path = _resolve_path_value(explicit, app_name, for_output=False) if explicit else ""

    open_value, open_action = _first_action_path(actions, _open_action_names(app_name))
    open_path = _resolve_path_value(open_value, app_name, for_output=False) if open_value else ""

    save_as_value, save_action = _first_action_path(
        actions,
        _save_as_action_names(app_name),
        path_keys=("filename", "path", "file_path", "output_path"),
    )
    save_as_path = _resolve_path_value(save_as_value, app_name, for_output=True) if save_as_value else ""

    named = _extract_named_file_path(command_text, app_name)
    named_path = _resolve_path_value(named, app_name, for_output=True) if named else ""
    fresh = _is_fresh_file_intent(app_name, command_text, actions)

    source_path = open_path or (explicit_path if explicit_path and not fresh else "")
    if source_path and not Path(source_path).exists():
        return {
            "success": False,
            "error_code": "FILE_NOT_FOUND",
            "message": f"Office input file was not found: {source_path}",
            "details": source_path,
            "app_type": app_name,
            "action_type": action_type,
        }

    if save_as_path:
        output_path = save_as_path
    elif explicit_path:
        output_path = explicit_path
    elif named_path:
        output_path = _next_available_path(named_path) if fresh else named_path
    elif source_path:
        output_path = source_path
    else:
        output_path = _generate_new_output_path(app_name)

    output = Path(output_path)
    try:
        output.parent.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        return {
            "success": False,
            "error_code": "INVALID_FILE_PATH",
            "message": f"Could not create output directory: {output.parent}",
            "details": str(exc),
            "app_type": app_name,
            "action_type": action_type,
        }

    if save_action is not None:
        save_action["filename"] = str(output.resolve())
    if open_action is not None and open_path:
        open_action["path"] = open_path

    return {
        "success": True,
        "app_type": app_name,
        "action_type": action_type,
        "source_path": source_path,
        "file_path": str(output.resolve()),
        "output_path": str(output.resolve()),
        "reason": (
            "frontend file path" if explicit_path else
            "open action path" if open_path else
            "save-as action path" if save_as_path else
            "named path from command" if named_path else
            "generated default output path"
        ),
    }


def _action_names(actions):
    return {
        str(action.get("action", "")).strip().lower()
        for action in (actions or [])
        if isinstance(action, dict)
    }


def _default_create_action(app_name, command_text):
    text = (command_text or "").lower()
    if not any(_contains_term(text, term) for term in OFFICE_ACTION_KEYWORDS):
        return None
    create_action = {
        "excel": "create_workbook",
        "word": "create_document",
        "powerpoint": "create_presentation",
        "ppt": "create_presentation",
    }.get(app_name)
    return {"action": create_action} if create_action else None


def _detect_action_type(command_text, actions):
    names = _action_names(actions)
    if names & {"create_workbook", "create_document", "create_presentation"}:
        return "create"
    if names & {"open_workbook", "open_document", "open_presentation"}:
        return "open"

    text = (command_text or "").lower()
    if any(_contains_term(text, term) for term in ("create", "make", "generate", "build", "new")):
        return "create"
    if any(_contains_term(text, term) for term in ("open", "load", "import")):
        return "open"
    if names or any(_contains_term(text, term) for term in OFFICE_ACTION_KEYWORDS):
        return "edit"
    return "unknown"


def _is_fresh_file_intent(app_name, command_text, actions):
    action_names = _action_names(actions)
    create_actions = {
        "excel": {"create_workbook"},
        "word": {"create_document"},
        "powerpoint": {"create_presentation"},
        "ppt": {"create_presentation"},
    }
    if action_names & create_actions.get(app_name, set()):
        return True

    text = (command_text or "").lower()
    creation_words = ("create", "new", "start", "make")
    target_words = ("file", "workbook", "document", "presentation", "ppt")
    return any(w in text for w in creation_words) and any(w in text for w in target_words)


def _should_start_fresh(app_name, command_text, actions, file_path):
    if file_path:
        return False
    if _extract_named_file_path(command_text, app_name):
        return False

    open_actions = {
        "excel": {"open_workbook"},
        "word": {"open_document"},
        "powerpoint": {"open_presentation"},
        "ppt": {"open_presentation"},
    }
    return not bool(_action_names(actions) & open_actions.get(app_name, set()))


def _ensure_fresh_file_action(app_name, command_text, actions, file_path):
    actions = list(actions or [])
    if not actions or not _should_start_fresh(app_name, command_text, actions, file_path):
        return actions

    create_action = {
        "excel": "create_workbook",
        "word": "create_document",
        "powerpoint": "create_presentation",
        "ppt": "create_presentation",
    }.get(app_name)
    if not create_action:
        return actions
    if str(actions[0].get("action", "")).strip().lower() == create_action:
        return actions

    logging.info(f"Prepending {create_action} for fresh {app_name} file: {command_text}")
    return [{"action": create_action}, *actions]


def _expand_powerpoint_slide_count(app_name, command_text, actions):
    if _canonical_office_app(app_name) != "powerpoint":
        return actions
    text = (command_text or "").lower()
    match = re.search(r"\b(?:create|make|generate|build|add)\b.*?\b(\d{1,2})\s+slides?\b", text)
    if not match:
        return actions
    target_count = max(1, min(int(match.group(1)), 50))
    existing = sum(
        1 for action in actions
        if isinstance(action, dict) and str(action.get("action", "")).lower() == "add_slide"
    )
    if existing >= target_count:
        return actions
    return list(actions) + [{"action": "add_slide", "layout": "title_content"} for _ in range(target_count - existing)]


def _resolve_output_file_path(app_name, command_text, actions, file_path):
    explicit = (file_path or "").strip()
    if explicit:
        return os.path.abspath(explicit)

    named = _extract_named_file_path(command_text, app_name)
    if named:
        # For "create/new file" style commands, avoid reusing locked/existing targets.
        if _is_fresh_file_intent(app_name, command_text, actions):
            return _next_available_path(named)
        return named

    if _should_start_fresh(app_name, command_text, actions, ""):
        return _generate_new_output_path(app_name)

    return ""


def _office_dependency_error(app_name):
    module_name, package_name = OFFICE_DEPENDENCIES.get(app_name, (None, None))
    if not module_name:
        return None
    try:
        __import__(module_name)
        return None
    except ModuleNotFoundError:
        return (
            f"{app_name.title()} support requires `{package_name}`. "
            f"Install it with `pip install {package_name}` or `pip install -r requirements.txt`."
        )


def _known_office_actions(app_name):
    cls = {
        "excel": ExcelExecutor,
        "word": WordExecutor,
        "powerpoint": PowerPointExecutor,
        "ppt": PowerPointExecutor,
    }.get(app_name)
    if not cls:
        return set()
    return {
        name[4:]
        for name in dir(cls)
        if name.startswith("_do_")
    }


def _has_explicit_save_action(app_name, actions, command_text="", file_path=""):
    names = _action_names(actions)
    save_map = {
        "excel": {"save_workbook", "save_workbook_as"},
        "word": {"save_document", "save_document_as"},
        "powerpoint": {"save_presentation", "save_presentation_as"},
        "ppt": {"save_presentation", "save_presentation_as"},
    }
    if names & save_map.get(app_name, set()):
        return True

    # Treat an explicit target filename/path as save intent.
    if (file_path or "").strip():
        return True
    if _extract_named_file_path(command_text, app_name):
        return True
    return False


def _run_office_actions(app_name, actions, file_path=None, command_text="", source_path=None):
    app_name = _canonical_office_app(app_name)
    output_path = str(Path((file_path or OFFICE_OUTPUTS.get(app_name, "output.xlsx"))).resolve())
    source_path = str(Path(source_path).resolve()) if source_path else ""
    executed = []
    failures = []
    opened = False
    persisted = False
    dependency_error = _office_dependency_error(app_name)

    logging.info(
        "Office execution start: app=%s source=%s output=%s actions=%s",
        app_name,
        source_path or "<new>",
        output_path,
        len(actions or []),
    )

    if dependency_error:
        failures.append(dependency_error)
        return {
            "success": False,
            "error_code": "OFFICE_DEPENDENCY_MISSING",
            "ok_count": 0,
            "total": len(actions or []),
            "executed": executed,
            "failures": failures,
            "output_path": output_path,
            "opened": opened,
            "persisted": persisted,
            "dependency_error": dependency_error,
        }

    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)

        if app_name == "excel":
            from openpyxl import Workbook, load_workbook
            load_path = source_path or (output_path if os.path.exists(output_path) else "")
            wb = load_workbook(load_path) if load_path else Workbook()
            ws = wb.active
            setattr(wb, "_path", output_path)
            executor = ExcelExecutor(wb, ws)

            for action in actions or []:
                current_wb = getattr(executor, "wb", wb)
                setattr(current_wb, "_path", output_path)
                ok = bool(executor.run(action))
                current_wb = getattr(executor, "wb", current_wb)
                setattr(current_wb, "_path", output_path)
                action_name = action.get("action", "unknown")
                if ok:
                    executed.append(action_name)
                else:
                    failures.append(f"{action_name} failed")

            final_obj = getattr(executor, "wb", wb)
            save_method = final_obj.save

        elif app_name == "word":
            from docx import Document
            load_path = source_path or (output_path if os.path.exists(output_path) else "")
            doc = Document(load_path) if load_path else Document()
            setattr(doc, "_path", output_path)
            executor = WordExecutor(doc)

            for action in actions or []:
                current_doc = getattr(executor, "doc", doc)
                setattr(current_doc, "_path", output_path)
                ok = bool(executor.run(action))
                current_doc = getattr(executor, "doc", current_doc)
                setattr(current_doc, "_path", output_path)
                action_name = action.get("action", "unknown")
                if ok:
                    executed.append(action_name)
                else:
                    failures.append(f"{action_name} failed")

            final_obj = getattr(executor, "doc", doc)
            save_method = final_obj.save

        elif app_name == "powerpoint":
            from pptx import Presentation
            load_path = source_path or (output_path if os.path.exists(output_path) else "")
            prs = Presentation(load_path) if load_path else Presentation()
            setattr(prs, "_path", output_path)
            executor = PowerPointExecutor(prs)

            for action in actions or []:
                current_prs = getattr(executor, "prs", prs)
                setattr(current_prs, "_path", output_path)
                ok = bool(executor.run(action))
                current_prs = getattr(executor, "prs", current_prs)
                setattr(current_prs, "_path", output_path)
                action_name = action.get("action", "unknown")
                if ok:
                    executed.append(action_name)
                else:
                    failures.append(f"{action_name} failed")

            final_obj = getattr(executor, "prs", prs)
            save_method = final_obj.save

        else:
            return {
                "success": False,
                "error_code": "UNSUPPORTED_OFFICE_APP",
                "ok_count": len(executed),
                "total": len(actions or []),
                "executed": executed,
                "failures": [f"Unsupported app: {app_name}"],
                "output_path": output_path,
                "persisted": False,
                "opened": False,
            }

        # Office automation persists successful actions by default. For an
        # existing source file, output_path is the same path unless save-as was
        # requested; for a new file, output_path is a generated outputs/ file.
        if not failures:
            try:
                save_method(output_path)
                persisted = Path(output_path).exists()
            except PermissionError:
                fallback_path = _next_available_path(output_path)
                save_method(fallback_path)
                output_path = fallback_path
                persisted = Path(output_path).exists()
                logging.warning("%s target was locked. Saved to fallback path: %s", app_name, output_path)
            except Exception as exc:
                failures.append(str(exc))
                logging.error("Office save failed for %s: %s", output_path, exc)

        if not failures and not persisted:
            failures.append("Save did not create a file on disk.")

    except Exception as exc:
        failures.append(str(exc))
        logging.error("Office execution failed: %s\n%s", exc, traceback.format_exc())

    if not failures and persisted and os.path.exists(output_path):
        try:
            opened = bool(system_core.open_path(output_path))
        except Exception:
            opened = False

    success = not failures and persisted and Path(output_path).exists()
    logging.info(
        "Office execution result: app=%s success=%s saved=%s output=%s",
        app_name,
        success,
        persisted,
        output_path,
    )

    return {
        "success": success,
        "error_code": "" if success else "OFFICE_SAVE_FAILED",
        "ok_count": len(executed),
        "total": len(actions or []),
        "executed": executed,
        "failures": failures,
        "output_path": output_path,
        "persisted": persisted,
        "opened": opened,
    }


def _handle_global_command(raw_text):
    """Handles system-wide agent: <app>: <command> triggers."""
    try:
        app_name, command = _extract_office_agent_command(raw_text)
        if app_name and command:
            if app_name == "ppt":
                app_name = "powerpoint"
            cache_key, actions, source, action_error = _resolve_actions(app_name, command)
            if action_error:
                logging.warning("Global office action parse failed: %s", action_error.message)
                return
            if not actions:
                logging.warning(f"No office action match for global command: {app_name}: {command}")
                return
            actions = _ensure_fresh_file_action(app_name, command, actions, "")
            actions = _expand_powerpoint_slide_count(app_name, command, actions)
            try:
                actions = validate_actions(app_name, actions, known_actions=_known_office_actions(app_name))
            except OfficeActionError as exc:
                logging.warning("Global office validation failed: %s", exc.message)
                return
            resolution = resolve_office_file_path({"raw": command}, actions, app_name)
            if not resolution.get("success"):
                logging.warning("Global office path resolution failed: %s", resolution.get("message"))
                return
            summary = _run_office_actions(
                app_name,
                actions,
                file_path=resolution["file_path"],
                source_path=resolution.get("source_path"),
                command_text=command,
            )
            if summary["failures"] and cache_key:
                command_map.remove_action(app_name, cache_key)
            logging.info(
                f"Global office [{source}] {app_name}: {command} -> "
                f"{summary['ok_count']}/{summary['total']} | {summary['output_path']}"
            )
            if summary.get("persisted"):
                _safe_speak(f"Executed {summary['ok_count']} actions in {app_name}")
            else:
                _safe_speak(f"Executed {summary['ok_count']} actions in {app_name}, not saved")
            return

        txt = (raw_text or "").strip()
        low = txt.lower()
        if low.startswith("agent "):
            sys_cmd = txt[len("agent "):].strip()
            sys_cmd = sys_cmd.replace("  ", " ").strip(" .,:;!?")
            if sys_cmd.startswith(("open ", "launch ", "start ", "run ", "boot ")):
                success, message = system_core.find_and_launch(sys_cmd)
                _safe_speak(
                    f"Opening {system_core.normalize_app_name(sys_cmd)}"
                    if success else f"Could not open {sys_cmd}"
                )
                logging.info(f"Voice system open [{sys_cmd}] => {success}: {message}")
            elif sys_cmd.startswith(("close ", "shut ", "exit ")):
                success, message = system_core.close_app(sys_cmd)
                _safe_speak(
                    f"Closing {system_core.normalize_app_name(sys_cmd)}"
                    if success else f"Could not close {sys_cmd}"
                )
                logging.info(f"Voice system close [{sys_cmd}] => {success}: {message}")
    except Exception as e:
        logging.error(f"Global command error: {e}\n{traceback.format_exc()}")


# Patch the keyboard/voice listener callbacks now that _handle_global_command is defined
_keyboard_listener.on_command = _handle_global_command
if _voice_listener:
    _voice_listener.on_command = _handle_global_command


# ===========================================================================
# ROUTES
# ===========================================================================

def _json_success(message, intent="unknown", **extra):
    payload = {
        "success": True,
        "status": "success",
        "intent": intent,
        "message": message,
    }
    payload.update(extra)
    return jsonify(payload)


def _json_error(message, intent="unknown", error_code="COMMAND_FAILED", http_status=200, **extra):
    payload = {
        "success": False,
        "status": "fail",
        "intent": intent,
        "error_code": error_code,
        "message": message,
    }
    payload.update(extra)
    return jsonify(payload), http_status


@app.route("/")
def index():
    html = render_template("index.html")
    patch_tag = '<script src="/static/reliability.js"></script>'
    if patch_tag not in html and "</body>" in html:
        html = html.replace("</body>", f"{patch_tag}</body>")
    return html


# ---- System commands ------------------------------------------------------

@app.route("/execute", methods=["POST"])
def execute():
    try:
        data = request.get_json(silent=True) or {}
        raw_cmd = (data.get("command") or "").strip()
        cmd = raw_cmd.lower()
        logging.info("Received command: %s", raw_cmd)

        office_route = _detect_office_intent(raw_cmd)
        if office_route.get("is_office"):
            logging.info(
                "Routing decision: original=%r intent=office_automation app=%s handler=office_execute reason=%s",
                raw_cmd,
                office_route.get("app_type"),
                office_route.get("reason"),
            )
            return _office_execute_impl({
                **data,
                "command": raw_cmd,
                "raw": office_route.get("command") or raw_cmd,
                "app": office_route.get("app_type"),
            })

        if cmd.startswith(("close ", "shut ", "exit ")):
            app_name = cmd.replace("close ", "").replace("shut ", "").replace("exit ", "").strip()
            success, message = system_core.close_app(app_name)
            _safe_speak(f"Closing {app_name}" if success else f"Could not close {app_name}")
            if success:
                return _json_success(message, intent="app_close", app_type=app_name)
            return _json_error(message, intent="app_close", error_code="APP_CLOSE_FAILED", app_type=app_name)

        if not _is_app_launch_command(cmd):
            logging.info(
                "Routing decision: original=%r intent=unknown handler=none reason=not an app-launch or office command",
                raw_cmd,
            )
            return _json_error(
                "Command was not recognized as an Office automation or app-launch request.",
                intent="unknown",
                error_code="UNKNOWN_COMMAND",
            )

        app_name = system_core.normalize_app_name(raw_cmd)
        logging.info(
            "Routing decision: original=%r intent=app_launch app=%s handler=system_core.find_and_launch reason=launch verb",
            raw_cmd,
            app_name,
        )
        success, message = system_core.find_and_launch(app_name)
        if success:
            _safe_speak(f"Opening {app_name}")
            return _json_success(message, intent="app_launch", app_type=app_name)

        if _is_known_office_app(app_name):
            logging.warning("Office app launch failed without manual selector: app=%s message=%s", app_name, message)
            return _json_error(
                f"Could not open configured Office application: {app_name}.",
                intent="app_launch",
                error_code="OFFICE_APP_LAUNCH_FAILED",
                app_type=app_name,
                details=message,
            )

        _safe_speak(f"I couldn't find {app_name}. Please select it manually.")
        path = ui.manual_selector()
        if path:
            norm_app = system_core.normalize_app_name(app_name)
            config.save_memory(norm_app, path, is_store_app=False)
            launched = system_core.open_path(path)
            if launched:
                _safe_speak("Path saved. Opening now.")
                return _json_success(
                    "Manual Selection Saved",
                    intent="app_launch",
                    app_type=norm_app,
                    file_path=path,
                    requires_manual_selection=True,
                )
            return _json_error(
                "Saved path, but launch failed",
                intent="app_launch",
                error_code="APP_LAUNCH_FAILED",
                app_type=norm_app,
                file_path=path,
                requires_manual_selection=True,
            )

        return _json_error(
            "Cancelled",
            intent="app_launch",
            error_code="MANUAL_SELECTION_CANCELLED",
            app_type=app_name,
            requires_manual_selection=True,
        )

    except Exception as e:
        logging.error("Command route error: %s\n%s", e, traceback.format_exc())
        return _json_error(
            "Command execution failed.",
            intent="unknown",
            error_code="COMMAND_ROUTE_ERROR",
            details=str(e),
            http_status=500,
        )


# ---- Office Agent ---------------------------------------------------------

@app.route("/office/execute", methods=["POST"])
def office_execute():
    try:
        return _office_execute_impl(request.get_json(silent=True) or {})
    except Exception as e:
        logging.error("Office route error: %s\n%s", e, traceback.format_exc())
        return _json_error(
            "Office command execution failed.",
            intent="office_automation",
            error_code="OFFICE_ROUTE_ERROR",
            details=str(e),
            http_status=500,
        )


def _office_execute_impl(data):
    app_name = (data.get("app") or "").lower().strip()
    command = (data.get("raw") or "").strip()
    full = (data.get("command") or "").strip()

    if not command and full:
        parsed_app, parsed_command = _extract_office_agent_command(full)
        if parsed_app and not app_name:
            app_name = (parsed_app or "").strip()
        if parsed_command:
            command = (parsed_command or "").strip()
        elif app_name:
            command = full

    if not app_name:
        office_route = _detect_office_intent(full or command)
        if office_route.get("is_office"):
            app_name = office_route.get("app_type", "")
            command = office_route.get("command") or command or full

    app_name = _canonical_office_app(app_name)

    if app_name not in OFFICE_APPS or not command:
        return _json_error(
            "Missing or invalid Office app/command.",
            intent="office_automation",
            error_code="INVALID_OFFICE_REQUEST",
            app_type=app_name or "unknown",
            action_type="unknown",
        )

    logging.info("Office request: original=%r app=%s", full or command, app_name)

    cache_key, actions, source, action_error = _resolve_actions(app_name, command)
    if action_error:
        logging.warning("Office action parse error: %s", action_error.message)
        return _json_error(
            action_error.message,
            intent="office_automation",
            error_code=action_error.error_code,
            app_type=app_name,
            action_type=_detect_action_type(command, actions),
            details=action_error.details,
        )
    if not actions:
        return _json_error(
            "No matching Office command found. Try a more specific action like 'create a new workbook' or 'add heading Introduction'.",
            intent="office_automation",
            error_code="NO_OFFICE_ACTION_MATCH",
            app_type=app_name,
            action_type="unknown",
            source=source,
        )

    requested_file_path = (data.get("file_path") or data.get("file") or "").strip()
    actions = _ensure_fresh_file_action(app_name, command, actions, requested_file_path)
    actions = _expand_powerpoint_slide_count(app_name, command, actions)
    try:
        actions = validate_actions(app_name, actions, known_actions=_known_office_actions(app_name))
    except OfficeActionError as exc:
        logging.warning("Office action validation failed: %s", exc.message)
        return _json_error(
            exc.message,
            intent="office_automation",
            error_code=exc.error_code,
            app_type=app_name,
            action_type=_detect_action_type(command, actions),
            details=exc.details,
            source=source,
        )

    resolution = resolve_office_file_path(data, actions, app_name, mode=_detect_action_type(command, actions))
    if not resolution.get("success"):
        logging.warning("Office path resolution failed: %s", resolution.get("message"))
        return _json_error(
            resolution.get("message", "Could not resolve Office file path."),
            intent="office_automation",
            error_code=resolution.get("error_code", "INVALID_FILE_PATH"),
            app_type=app_name,
            action_type=resolution.get("action_type", "unknown"),
            details=resolution.get("details", ""),
            source=source,
        )

    logging.info(
        "Routing decision: original=%r intent=office_automation app=%s action_type=%s handler=office_executor reason=%s output=%s action_count=%s",
        command,
        app_name,
        resolution.get("action_type"),
        resolution.get("reason"),
        resolution.get("file_path"),
        len(actions),
    )

    summary = _run_office_actions(
        app_name,
        actions,
        file_path=resolution["file_path"],
        source_path=resolution.get("source_path"),
        command_text=command,
    )
    if summary.get("dependency_error"):
        return _json_error(
            summary["dependency_error"],
            intent="office_automation",
            error_code="OFFICE_DEPENDENCY_MISSING",
            app_type=app_name,
            action_type=resolution.get("action_type", "unknown"),
            source=source,
            file_path=summary["output_path"],
            output_file=summary["output_path"],
        )
    if summary["failures"] and cache_key:
        command_map.remove_action(app_name, cache_key)
    if summary["failures"]:
        return _json_error(
            f"Could not save {app_name.title()} file.",
            intent="office_automation",
            error_code=summary.get("error_code") or "OFFICE_SAVE_FAILED",
            app_type=app_name,
            action_type=resolution.get("action_type", "unknown"),
            details=f"{summary['ok_count']}/{summary['total']} done | {' | '.join(summary['failures'])}",
            source=source,
            file_path=summary["output_path"],
            output_file=summary["output_path"],
            persisted=summary.get("persisted", False),
        )

    app_label = {"excel": "Excel", "word": "Word", "powerpoint": "PowerPoint"}.get(app_name, app_name.title())
    return _json_success(
        f"Created {app_label} file successfully." if resolution.get("action_type") == "create" else f"Updated {app_label} file successfully.",
        intent="office_automation",
        app_type=app_name,
        action_type=resolution.get("action_type", "unknown"),
        file_path=summary["output_path"],
        source=source,
        output_file=summary["output_path"],
        persisted=summary.get("persisted", False),
        opened=summary.get("opened", False),
        action_count=summary.get("total", len(actions)),
        executed=summary.get("executed", []),
    )


@app.route("/command", methods=["POST"])
def office_command():
    try:
        return _office_execute_impl(request.get_json(silent=True) or {})
    except Exception as e:
        logging.error("Command office route error: %s\n%s", e, traceback.format_exc())
        return _json_error(
            "Office command execution failed.",
            intent="office_automation",
            error_code="OFFICE_ROUTE_ERROR",
            details=str(e),
            http_status=500,
        )


# ---- Voice control --------------------------------------------------------

@app.route("/voice/status", methods=["GET"])
def voice_status():
    if not _voice_listener:
        return jsonify(
            status="fail", available=False, enabled=False,
            message="Voice module unavailable. Install SpeechRecognition + PyAudio."
        )
    heard = _voice_listener.last_heard
    if time.time() - (_voice_listener.last_heard_at or 0) > 8:
        heard = ""
    return jsonify(
        status="success",
        available=_voice_listener.available,
        enabled=_voice_listener.is_running,
        armed=_voice_listener.armed,
        armed_seconds=round(_voice_listener.armed_seconds_left, 1),
        heard=heard,
        error=_voice_listener.last_error
    )


@app.route("/voice/start", methods=["POST"])
def voice_start():
    if not _voice_listener:
        return jsonify(status="fail", message="Voice module unavailable")
    ok = _voice_listener.start()
    voice_state["enabled"] = bool(ok)
    return jsonify(
        status="success" if ok else "fail",
        message="Voice listener started" if ok else (_voice_listener.last_error or "Could not start voice listener")
    )


@app.route("/voice/stop", methods=["POST"])
def voice_stop():
    if not _voice_listener:
        return jsonify(status="fail", message="Voice module unavailable")
    _voice_listener.stop()
    voice_state["enabled"] = False
    return jsonify(status="success", message="Voice listener stopped")


# ---- OCR ------------------------------------------------------------------

@app.route("/ocr/snip", methods=["POST"])
def ocr_snip():
    try:
        if not OCR_AVAILABLE:
            return jsonify(status="fail", message="OCR not available")
        ocr_utils.snip_queue.put("snip")
        try:
            path = ocr_utils.result_queue.get(timeout=60)
        except Exception:
            return jsonify(status="fail", message="Snip timed out")
        if not path:
            return jsonify(status="fail", message="Snip cancelled")
        text = ocr_utils.image_to_text(path)
        last_ocr["text"] = text
        last_ocr["pending"] = False
        return jsonify(status="success", text=text)
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/ocr/screenshot", methods=["POST"])
def ocr_screenshot():
    try:
        if not OCR_AVAILABLE:
            return jsonify(status="fail", message="OCR not available")
        path = ocr_utils.capture_fullscreen()
        text = ocr_utils.image_to_text(path)
        last_ocr["text"] = text
        last_ocr["pending"] = False
        return jsonify(status="success", text=text)
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/ocr/file", methods=["POST"])
def ocr_file():
    try:
        if not OCR_AVAILABLE:
            return jsonify(status="fail", message="OCR not available")
        path = ui.file_selector(
            "Select an Image File",
            [("Images", "*.png *.jpg *.jpeg *.bmp *.tiff"), ("All Files", "*.*")]
        )
        if not path:
            return jsonify(status="fail", message="No file selected")
        text = ocr_utils.image_to_text(path)
        last_ocr["text"] = text
        last_ocr["pending"] = False
        return jsonify(status="success", text=text, message=f"OCR complete — {len(text)} chars")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/ocr/read", methods=["POST"])
def ocr_read():
    try:
        text = last_ocr.get("text", "")
        if not text:
            return jsonify(status="fail", message="No OCR text. Run OCR first.")
        threading.Thread(target=ocr_utils.speak_text, args=(text,), daemon=True).start()
        return jsonify(status="success", message="Speaking...")
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/ocr/stop_read", methods=["POST"])
def ocr_stop_read():
    try:
        ocr_utils.stop_speaking()
        return jsonify(status="success", message="Stopped")
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/ocr/poll", methods=["GET"])
def ocr_poll():
    try:
        if last_ocr.get("pending"):
            last_ocr["pending"] = False
            return jsonify(
                status="ready",
                text=last_ocr["text"],
                message=f"Hotkey OCR complete — {len(last_ocr['text'])} chars"
            )
        return jsonify(status="waiting")
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/ocr/save_txt", methods=["POST"])
def ocr_save_txt():
    try:
        text = last_ocr.get("text", "")
        if not text:
            return jsonify(status="fail", message="No OCR text. Run OCR first.")
        path = ocr_utils.save_as_txt(text)
        if not path:
            return jsonify(status="fail", message="Save cancelled.")
        return jsonify(status="success", message=f"Saved: {path}")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/ocr/save_pdf", methods=["POST"])
def ocr_save_pdf():
    try:
        text = last_ocr.get("text", "")
        if not text:
            return jsonify(status="fail", message="No OCR text. Run OCR first.")
        if not PDF_AVAILABLE:
            return jsonify(status="fail", message="Install fpdf2: pip install fpdf2")
        path = pdf_utils.create_report(text, title="OCR Result")
        if not path:
            return jsonify(status="fail", message="Save cancelled.")
        return jsonify(status="success", message=f"Saved: {path}")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/ocr/clipboard", methods=["POST"])
def ocr_clipboard():
    try:
        text = last_ocr.get("text", "")
        if not text:
            return jsonify(status="fail", message="No OCR text. Run OCR first.")
        ocr_utils.copy_to_clipboard(text)
        return jsonify(status="success", message="Copied to clipboard")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


# ---- PDF Reader -----------------------------------------------------------

@app.route("/reader/open", methods=["POST"])
def reader_open():
    try:
        if not READER_AVAILABLE:
            return jsonify(status="fail", message="PDF reader module not found")
        path = ui.file_selector("Select PDF to Read", [("PDFs", "*.pdf")])
        if not path:
            return jsonify(status="fail", message="No file selected")
        threading.Thread(target=pdf_reader.start_reading, args=(path, 0), daemon=True).start()
        time.sleep(0.5)
        return jsonify(status="success", message="Reading started", **pdf_reader.get_status())
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/reader/pause", methods=["POST"])
def reader_pause():
    try:
        pdf_reader.pause_reading()
        return jsonify(status="success", message="Paused", **pdf_reader.get_status())
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/reader/resume", methods=["POST"])
def reader_resume():
    try:
        pdf_reader.resume_reading()
        return jsonify(status="success", message="Resumed", **pdf_reader.get_status())
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/reader/stop", methods=["POST"])
def reader_stop():
    try:
        pdf_reader.stop_reading()
        return jsonify(status="success", message="Stopped")
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/reader/next", methods=["POST"])
def reader_next():
    try:
        pdf_reader.next_page()
        return jsonify(status="success", **pdf_reader.get_status())
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/reader/prev", methods=["POST"])
def reader_prev():
    try:
        pdf_reader.prev_page()
        return jsonify(status="success", **pdf_reader.get_status())
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/reader/speed", methods=["POST"])
def reader_speed():
    try:
        data = request.json
        pdf_reader.set_speed(data.get("speed", 150))
        return jsonify(status="success", message=f"Speed: {data.get('speed')} WPM")
    except Exception as e:
        return jsonify(status="fail", message=str(e))


@app.route("/reader/status", methods=["GET"])
def reader_status():
    try:
        return jsonify(pdf_reader.get_status())
    except Exception:
        return jsonify(is_reading=False, is_paused=False, current_page=0, total_pages=0, speed=150)


# ---- PDF Tools ------------------------------------------------------------

@app.route("/pdf/merge", methods=["POST"])
def pdf_merge():
    try:
        if not PDF_AVAILABLE:
            return jsonify(status="fail", message="Install pypdf: pip install pypdf")
        paths = pdf_utils.ask(
            kind="openmultiple",
            title="Select PDFs to Merge (hold Ctrl for multiple)",
            filetypes=[("PDF Files", "*.pdf"), ("All Files", "*.*")]
        )
        if not paths:
            return jsonify(status="fail", message="No files selected.")
        out = pdf_utils.merge_pdfs(paths)
        if not out:
            return jsonify(status="fail", message="Save cancelled.")
        return jsonify(status="success", message=f"Merged {len(paths)} PDFs → {out}")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/pdf/split", methods=["POST"])
def pdf_split():
    try:
        if not PDF_AVAILABLE:
            return jsonify(status="fail", message="Install pypdf: pip install pypdf")
        path = ui.file_selector("Select PDF to Split", [("PDFs", "*.pdf")])
        if not path:
            return jsonify(status="fail", message="No file selected.")
        pages = pdf_utils.split_pdf(path)
        if not pages:
            return jsonify(status="fail", message="Save cancelled or no pages.")
        return jsonify(status="success", message=f"Split into {len(pages)} files")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/pdf/create", methods=["POST"])
def pdf_create():
    try:
        if not PDF_AVAILABLE:
            return jsonify(status="fail", message="Install fpdf2: pip install fpdf2")
        data = request.json
        text = data.get("text", "").strip()
        title = (data.get("title", "Report") or "Report").strip()
        if not text:
            return jsonify(status="fail", message="No text provided")
        path = pdf_utils.create_report(text, title=title)
        if not path:
            return jsonify(status="fail", message="Save cancelled.")
        return jsonify(status="success", message=f"PDF saved: {path}")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


# ---- PDF Editor -----------------------------------------------------------

@app.route("/editor/open", methods=["POST"])
def editor_open():
    try:
        if not PDF_EDITOR_AVAILABLE:
            return jsonify(status="fail", message="PDF Editor not available")
        path = ui.file_selector("Select PDF to Edit", [("PDF Files", "*.pdf")])
        if not path:
            return jsonify(status="fail", message="No file selected")
        data = pdf_editor.extract_pdf_text(path)
        if data.get("status") != "success":
            return jsonify(status="fail", message=data.get("message", "Failed to open PDF"))
        return jsonify(status="success", file_path=path, pages=data["pages"], total_pages=data["total_pages"])
    except Exception as e:
        logging.error(traceback.format_exc())
        return jsonify(status="fail", message=str(e))


@app.route("/editor/render-page", methods=["POST"])
def editor_render_page():
    try:
        if not PDF_EDITOR_AVAILABLE:
            return jsonify(status="fail", message="PDF Editor not available")
        data = request.json
        pdf_path = data.get("file_path")
        page_num = data.get("page_num", 0)
        if not pdf_path:
            return jsonify(status="fail", message="No file path provided")
        result = pdf_editor.render_page_as_image(pdf_path, page_num)
        if result.get("status") != "success":
            return jsonify(status="fail", message=result.get("message", "Render failed"))
        return jsonify(status="success", **{k: v for k, v in result.items() if k != "status"})
    except Exception as e:
        logging.error(traceback.format_exc())
        return jsonify(status="fail", message=str(e))


@app.route("/editor/save", methods=["POST"])
def editor_save():
    try:
        if not PDF_EDITOR_AVAILABLE:
            return jsonify(status="fail", message="PDF Editor not available")
        data = request.json
        pdf_path = data.get("file_path")
        edits = data.get("edits", [])
        if not pdf_path:
            return jsonify(status="fail", message="No file path provided")
        result = pdf_editor.save_edited_pdf(pdf_path, edits)
        if result.get("status") != "success":
            return jsonify(status="fail", message=result.get("message", "Save failed"))
        return jsonify(status="success", message=result.get("message", "Saved successfully"))
    except Exception as e:
        logging.error(traceback.format_exc())
        return jsonify(status="fail", message=str(e))


@app.route("/editor/detect-form", methods=["POST"])
def editor_detect_form():
    try:
        if not PDF_EDITOR_AVAILABLE:
            return jsonify(status="fail", message="PDF Editor not available")
        path = ui.file_selector("Select PDF", [("PDF Files", "*.pdf")])
        if not path:
            return jsonify(status="fail", message="No file selected")
        fields = pdf_editor.detect_form_fields(path)
        return jsonify(
            status="success",
            is_form=len(fields) > 0,
            field_count=len(fields),
            fields=list(fields.keys()),
            file_path=path
        )
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/editor/fill-form", methods=["POST"])
def editor_fill_form():
    try:
        if not PDF_EDITOR_AVAILABLE:
            return jsonify(status="fail", message="PDF Editor not available")
        data = request.json
        pdf_path = data.get("file_path")
        form_data = data.get("form_data", {})
        if not pdf_path:
            return jsonify(status="fail", message="No file path provided")
        result = pdf_editor.fill_form(pdf_path, form_data)
        if result:
            return jsonify(status="success", message="Form saved successfully")
        return jsonify(status="fail", message="Save cancelled")
    except Exception as e:
        print(traceback.format_exc())
        return jsonify(status="fail", message=f"Error: {str(e)}")


@app.route("/editor/get-field-options", methods=["POST"])
def editor_get_field_options():
    try:
        if not PDF_EDITOR_AVAILABLE:
            return jsonify(status="fail", message="PDF Editor not available")
        data = request.json
        pdf_path = data.get("file_path")
        field_name = data.get("field_name")
        options = pdf_editor.get_form_field_options(pdf_path, field_name)
        return jsonify(status="success", field_name=field_name, options=options)
    except Exception as e:
        return jsonify(status="fail", message=f"Error: {str(e)}")


# ===========================================================================
# ENTRY POINT
# ===========================================================================

if __name__ == "__main__":

    # ---- OCR snip overlay (must be on main thread) ------------------------
    if OCR_AVAILABLE:
        threading.Thread(target=ocr_utils.run_snip_overlay_main_thread, daemon=True).start()

    # ---- OCR keyboard hotkeys ---------------------------------------------
    if KEYBOARD_AVAILABLE and OCR_AVAILABLE:
        keyboard.add_hotkey(
            "ctrl+shift+s",
            lambda: threading.Thread(
                target=ocr_utils.trigger_snip_and_ocr, args=(last_ocr,), daemon=True
            ).start()
        )
        keyboard.add_hotkey(
            "ctrl+shift+f",
            lambda: threading.Thread(
                target=ocr_utils.trigger_screenshot_and_ocr, args=(last_ocr,), daemon=True
            ).start()
        )
        print("🔑  Ctrl+Shift+S → Snip OCR  |  Ctrl+Shift+F → Fullscreen OCR")

    # ---- Global Office Agent listeners ------------------------------------
    threading.Thread(target=_clipboard_listener.start, daemon=True, name="ClipboardListener").start()
    threading.Thread(target=_keyboard_listener.start, daemon=True, name="KeyboardListener").start()
    print("⌨️   Global agent listener active")
    print("     Type  agent: excel: <command>  anywhere + Enter")

    if _voice_listener and _voice_listener.available:
        if _voice_listener.start():
            voice_state["enabled"] = True
            print("Voice wake listener active (say: agent <app> <command>)")
        else:
            print(f"Voice listener not started: {_voice_listener.last_error}")

    # ---- Start Flask ------------------------------------------------------
    flask_thread = threading.Thread(
        target=lambda: app.run(host="127.0.0.1", port=5000, debug=False),
        daemon=True
    )
    flask_thread.start()
    time.sleep(1)

    # ---- Open browser -----------------------------------------------------
    webbrowser.open("http://127.0.0.1:5000")
    print("✅  Agent running at http://127.0.0.1:5000")

    # ---- Dialog listener must be on main thread ---------------------------
    if PDF_AVAILABLE:
        pdf_utils.run_dialog_listener()
    else:
        try:
            while True:
                time.sleep(1)
        except KeyboardInterrupt:
            print("\n👋 Agent stopped.")
