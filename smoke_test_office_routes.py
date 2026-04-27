from pathlib import Path
import json
import re

import server
from ai.openai_handler import OpenAIHandler
from executor.excel_executor import _normalize_excel_argb
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from docx import Document
from utils import command_map
from utils.app_alias_guard import validate_manual_app_alias
from utils.file_paths import resolve_existing_office_path
from utils.office_actions import OfficeActionError, normalize_actions, validate_actions


COMMAND_MAP = Path("command_map.json")
COMMAND_JSON_FILES = [
    Path("excel_commands.json"),
    Path("word_commands.json"),
    Path("powerpoint_commands.json"),
]
KNOWN_COMMAND_PLACEHOLDERS = {
    "after", "alignment", "angle", "animation_type", "audio_path", "base",
    "before", "bookmark_name", "bottom", "cell", "cell1", "cell2",
    "character", "chart_type", "col_index", "color", "color1", "color2",
    "cols", "column", "compare_path", "cond1", "cond2", "condition",
    "condition1", "condition2", "count", "criteria", "criteria1",
    "criteria2", "data_source", "date_format", "day", "decimals", "delay",
    "delimiter", "dest_cell", "direction", "end_date", "end_row",
    "error_value", "exponent", "false_val", "file_path", "filename",
    "find_text", "font_name", "footer_text", "format", "formula",
    "from_index", "header_text", "height", "icon_name", "image_path",
    "indent", "items", "layout", "left", "level", "lookup_array",
    "lookup_range", "lookup_value", "month", "new_name", "nper",
    "num_chars", "old_name", "order", "orientation", "output_path",
    "part", "password", "position", "pv", "range", "range1", "range2",
    "range_name", "rate", "replace_text", "result_cell", "return_array",
    "return_range", "right", "row_index", "row_number", "rows", "scheme",
    "seconds", "shape_type", "sheet_name", "size", "slide_index",
    "source_range", "spacing", "speed", "start", "start_cell",
    "start_date", "start_row", "style", "sum_range", "sumifs_formula",
    "table_range", "target", "target_cell", "text", "theme", "to_index",
    "top", "total", "transition_type", "true_val", "type", "unit", "url",
    "val1", "val2", "value", "values", "video_path", "width", "year",
    "zoom_level",
}


def _payload(response):
    data = response.get_json(silent=True)
    assert isinstance(data, dict), response.data
    assert "success" in data, data
    assert "status" in data, data
    assert "intent" in data, data
    assert "message" in data, data
    assert "data" in data, data
    return data


def _assert_office_file(data, suffix):
    assert data["success"] is True, data
    assert data["intent"] == "office_automation", data
    assert data.get("file_path"), data
    path = Path(data["file_path"])
    assert path.exists(), data
    assert path.suffix.lower() == suffix, data
    assert path.stat().st_size > 0, data
    assert data.get("persisted") is True, data
    return path


def _post_office(client, app, raw, **extra):
    return _payload(client.post("/office/execute", json={"app": app, "raw": raw, **extra}))


def _assert_no_polluted_known_apps():
    known_apps = Path("known_apps.json").read_text(encoding="utf-8").lower()
    assert "create new word file named hello" not in known_apps
    assert "word_testdocx" not in known_apps


def _audit_command_json():
    for path in COMMAND_JSON_FILES:
        commands = json.loads(path.read_text(encoding="utf-8"))
        app = path.stem.replace("_commands", "")
        if app == "powerpoint":
            app = "powerpoint"
        known_executor_actions = server._known_office_actions(app)
        assert isinstance(commands, dict) and commands, path
        for command_name, command in commands.items():
            params = command.get("params") or {}
            assert isinstance(params, dict), (path, command_name)
            assert isinstance(command.get("keywords", []), list), (path, command_name)
            action_name = command.get("action")
            assert action_name, (path, command_name)
            assert action_name in known_executor_actions, (path, command_name, action_name)
            for key, value in params.items():
                if not isinstance(value, str):
                    continue
                placeholders = re.findall(r"\{([^{}]+)\}", value)
                for placeholder in placeholders:
                    assert placeholder in KNOWN_COMMAND_PLACEHOLDERS, (path, command_name, key, placeholder)


def _audit_command_cache():
    if not COMMAND_MAP.exists():
        return
    data = json.loads(COMMAND_MAP.read_text(encoding="utf-8"))
    for app, command_map in data.items():
        known = server._known_office_actions(app)
        for command_text, actions in command_map.items():
            validate_actions(app, actions, known_actions=known)


def _assert_excel_fill(path, cell, expected):
    wb = load_workbook(path)
    rgb = wb.active[cell].fill.fgColor.rgb
    assert rgb == expected, f"{cell} fill was {rgb}, expected {expected}"


def main():
    original_command_map = COMMAND_MAP.read_bytes() if COMMAND_MAP.exists() else None

    # Do not launch real desktop apps or open blocking file pickers in smoke tests.
    server.system_core.open_path = lambda *args, **kwargs: True
    server.ui.manual_selector = lambda: ""

    try:
        _assert_no_polluted_known_apps()
        _audit_command_json()
        _audit_command_cache()

        ok, code, _ = validate_manual_app_alias("create new word file named hello")
        assert ok is False and code == "MANUAL_APP_ALIAS_REJECTED_DOCUMENT_COMMAND"
        assert validate_manual_app_alias("chrome")[0] is True

        assert _normalize_excel_argb("FF0000") == "FFFF0000"
        assert _normalize_excel_argb("#FF0000") == "FFFF0000"
        assert _normalize_excel_argb("yellow") == "FFFFFF00"
        try:
            validate_actions("excel", [{"action": "set_bg_color", "range": "WIDTHS:WIDTHS", "color": "yellow"}], known_actions=server._known_office_actions("excel"))
            raise AssertionError("Pseudo-ranges must be rejected")
        except OfficeActionError as exc:
            assert exc.error_code == "INVALID_EXCEL_RANGE"
        try:
            validate_actions("excel", [{"action": "set_bg_color", "range": "A1", "color": "not-a-color"}], known_actions=server._known_office_actions("excel"))
            raise AssertionError("Invalid colors must be rejected")
        except OfficeActionError as exc:
            assert exc.error_code == "INVALID_COLOR"
        try:
            validate_actions("excel", [{"action": "write_cell", "cell": "A1", "value": "{missing_placeholder}"}], known_actions=server._known_office_actions("excel"))
            raise AssertionError("Unresolved placeholders must be rejected")
        except OfficeActionError as exc:
            assert exc.error_code == "INVALID_ACTION"

        client = server.app.test_client()

        routing_cases = [
            ("create a new Excel file", ".xlsx"),
            ("create a new Word document", ".docx"),
            ("create a new PowerPoint presentation", ".pptx"),
            ("make a spreadsheet with 3 columns and 5 rows", ".xlsx"),
            ("create a presentation with 3 slides about sales performance", ".pptx"),
        ]

        generated = {}
        for command, suffix in routing_cases:
            data = _payload(client.post("/execute", json={"command": command}))
            generated[command] = _assert_office_file(data, suffix)
            assert not data.get("requires_manual_selection"), data

        three_slide_path = generated["create a presentation with 3 slides about sales performance"]
        assert len(Presentation(str(three_slide_path)).slides) == 3

        excel_table = _post_office(
            client,
            "excel",
            "create workbook and add table with 5 rows and 3 columns",
        )
        excel_table_path = _assert_office_file(excel_table, ".xlsx")
        wb = load_workbook(excel_table_path)
        assert len(wb.active.tables) >= 1
        assert wb.active["A1"].value == "Column1"

        excel_format = _post_office(
            client,
            "excel",
            "create workbook, add table with 5 rows and 3 columns, make header bold, set background color to yellow",
        )
        excel_format_path = _assert_office_file(excel_format, ".xlsx")
        wb = load_workbook(excel_format_path)
        assert wb.active["A1"].font.bold is True
        _assert_excel_fill(excel_format_path, "A1", "FFFFFF00")

        excel_contextual = _post_office(
            client,
            "excel",
            "create workbook; add table with 5 rows and 3 columns at A1; make the header bold; set the header background yellow; apply borders to the full table; autofit columns",
        )
        contextual_path = _assert_office_file(excel_contextual, ".xlsx")
        contextual_ws = load_workbook(contextual_path).active
        assert contextual_ws["A1"].font.bold is True
        assert contextual_ws["A1"].border.left.style == "thin"
        assert excel_contextual.get("plan", {}).get("context", {}).get("header_range") == "A1:C1", excel_contextual
        assert "WIDTHS:WIDTHS" not in json.dumps(excel_contextual), excel_contextual

        multiline = _post_office(
            client,
            "excel",
            "create workbook\nadd table with 3 columns and 5 rows\nmake header bold",
        )
        multiline_path = _assert_office_file(multiline, ".xlsx")
        assert load_workbook(multiline_path).active["A1"].font.bold is True

        values = _post_office(
            client,
            "excel",
            "create workbook; write Profit and Loss in A1; make A1 bold",
        )
        values_path = _assert_office_file(values, ".xlsx")
        ws = load_workbook(values_path).active
        assert ws["A1"].value == "Profit and Loss"
        assert ws["A1"].font.bold is True

        formula = _post_office(
            client,
            "excel",
            "create workbook. write formula in D7 that calculates the total from D3 to D6",
        )
        formula_path = _assert_office_file(formula, ".xlsx")
        assert load_workbook(formula_path, data_only=False).active["D7"].value == "=SUM(D3:D6)"

        concat = _post_office(client, "excel", "create workbook; concatenate A1 and B1 in C1")
        concat_path = _assert_office_file(concat, ".xlsx")
        assert load_workbook(concat_path, data_only=False).active["C1"].value == "=CONCATENATE(A1,B1)"

        textjoin_comma = _post_office(client, "excel", "create workbook; textjoin A1:A5 with comma into B1")
        textjoin_comma_path = _assert_office_file(textjoin_comma, ".xlsx")
        assert load_workbook(textjoin_comma_path, data_only=False).active["B1"].value == '=TEXTJOIN(",",TRUE,A1:A5)'

        textjoin_space = _post_office(client, "excel", "create workbook; textjoin A1:A5 with space into B1")
        textjoin_space_path = _assert_office_file(textjoin_space, ".xlsx")
        assert load_workbook(textjoin_space_path, data_only=False).active["B1"].value == '=TEXTJOIN(" ",TRUE,A1:A5)'

        sumifs_one = _post_office(client, "excel", "create workbook; sumifs C:C where A:A is North in D1")
        sumifs_one_path = _assert_office_file(sumifs_one, ".xlsx")
        assert load_workbook(sumifs_one_path, data_only=False).active["D1"].value == '=SUMIFS(C:C,A:A,"North")'

        sumifs_two = _post_office(client, "excel", "create workbook; sumifs C:C where A:A is North and B:B is Q1 in E1")
        sumifs_two_path = _assert_office_file(sumifs_two, ".xlsx")
        assert load_workbook(sumifs_two_path, data_only=False).active["E1"].value == '=SUMIFS(C:C,A:A,"North",B:B,"Q1")'

        today = _post_office(client, "excel", "create workbook; put today's date in A1")
        today_path = _assert_office_file(today, ".xlsx")
        assert load_workbook(today_path, data_only=False).active["A1"].value == "=TODAY()"

        now = _post_office(client, "excel", "create workbook; put current time in B1")
        now_path = _assert_office_file(now, ".xlsx")
        assert load_workbook(now_path, data_only=False).active["B1"].value == "=NOW()"

        time_short = _post_office(client, "excel", "create workbook; format time as hh:mm on A1")
        time_short_path = _assert_office_file(time_short, ".xlsx")
        assert load_workbook(time_short_path).active["A1"].number_format == "hh:mm"

        time_long = _post_office(client, "excel", "create workbook; format time as hh:mm:ss on A1")
        time_long_path = _assert_office_file(time_long, ".xlsx")
        assert load_workbook(time_long_path).active["A1"].number_format == "hh:mm:ss"

        protected = _post_office(client, "excel", "create workbook; protect active sheet with password secret")
        protected_path = _assert_office_file(protected, ".xlsx")
        assert load_workbook(protected_path).active.protection.sheet is True

        named_protected = _post_office(client, "excel", "create workbook; rename sheet to Sales; protect sheet named Sales with password secret")
        named_protected_path = _assert_office_file(named_protected, ".xlsx")
        assert load_workbook(named_protected_path)["Sales"].protection.sheet is True

        unprotected = _post_office(client, "excel", "unprotect sheet", file_path=str(protected_path))
        assert unprotected["success"] is True, unprotected
        assert load_workbook(protected_path).active.protection.sheet is False

        protected_again = _post_office(client, "excel", "protect active sheet with password secret", file_path=str(protected_path))
        assert protected_again["success"] is True, protected_again
        unprotected_password = _post_office(client, "excel", "unprotect sheet with password secret", file_path=str(protected_path))
        assert unprotected_password["success"] is True, unprotected_password
        assert load_workbook(protected_path).active.protection.sheet is False

        workbook_protected = _post_office(client, "excel", "protect workbook with password secret", file_path=str(protected_path))
        assert workbook_protected["success"] is True, workbook_protected
        workbook_unprotected = _post_office(client, "excel", "unprotect workbook with password secret", file_path=str(protected_path))
        assert workbook_unprotected["success"] is True, workbook_unprotected

        temp_input = Path("outputs/office/excel/open_action_input.xlsx").resolve()
        temp_input.parent.mkdir(parents=True, exist_ok=True)
        input_wb = Workbook()
        input_wb.active["B1"] = "Existing"
        input_wb.save(temp_input)
        opened = _post_office(
            client,
            "excel",
            f'open workbook "{temp_input}" and write Total in A1',
        )
        _assert_office_file(opened, ".xlsx")
        assert Path(opened["file_path"]).resolve() == temp_input
        assert load_workbook(temp_input).active["A1"].value == "Total"
        resolved_no_ext = resolve_existing_office_path("open_action_input", "excel", base_dir=temp_input.parent, command_text="open workbook open_action_input")
        assert resolved_no_ext.resolve() == temp_input
        opened_named = _post_office(
            client,
            "excel",
            "open workbook open_action_input and write Named in A2",
        )
        assert opened_named["success"] is True, opened_named
        assert load_workbook(temp_input).active["A2"].value == "Named"

        word = _post_office(
            client,
            "word",
            "create document. Add heading Summary. Add paragraph This is the first draft.",
        )
        word_path = _assert_office_file(word, ".docx")
        doc_text = "\n".join(p.text for p in Document(str(word_path)).paragraphs)
        assert "Summary" in doc_text
        assert "This is the first draft" in doc_text

        word_content = _post_office(
            client,
            "word",
            "create a Word document with heading Project Update and paragraph The project is on track",
        )
        word_content_path = _assert_office_file(word_content, ".docx")
        content_text = "\n".join(p.text for p in Document(str(word_content_path)).paragraphs)
        assert "Project Update" in content_text
        assert "The project is on track" in content_text

        existing_doc = Path("outputs/office/word/open_action_input.docx").resolve()
        existing_doc.parent.mkdir(parents=True, exist_ok=True)
        seed_doc = Document()
        seed_doc.add_paragraph("Existing")
        seed_doc.save(existing_doc)
        opened_doc = _post_office(
            client,
            "word",
            f'open document "{existing_doc}"; add paragraph Loaded document',
        )
        _assert_office_file(opened_doc, ".docx")
        opened_text = "\n".join(p.text for p in Document(str(existing_doc)).paragraphs)
        assert "Loaded document" in opened_text

        missing_doc = _post_office(client, "word", "open existing document missing_file_for_test.docx")
        assert missing_doc["success"] is False
        assert missing_doc["error_code"] == "FILE_NOT_FOUND"

        saved_doc = _post_office(client, "word", "create a Word document; save document")
        _assert_office_file(saved_doc, ".docx")

        save_as_doc = _post_office(client, "word", "create a Word document; save document as report_for_test")
        save_as_path = _assert_office_file(save_as_doc, ".docx")
        assert save_as_path.name.startswith("report_for_test")

        close_doc = _payload(client.post("/execute", json={"command": "close document"}))
        assert close_doc["intent"] == "office_automation", close_doc
        assert "close_document" in close_doc.get("executed", []), close_doc

        deck = _post_office(
            client,
            "powerpoint",
            "create presentation, add title slide, add agenda slide, add conclusion slide",
        )
        deck_path = _assert_office_file(deck, ".pptx")
        prs = Presentation(str(deck_path))
        assert len(prs.slides) == 3
        deck_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        assert "Agenda" in deck_text
        assert "Conclusion" in deck_text

        indexed = _post_office(client, "powerpoint", "create a presentation with 2 slides")
        indexed_path = _assert_office_file(indexed, ".pptx")
        first = _post_office(
            client,
            "powerpoint",
            "write First in title on slide 1",
            file_path=str(indexed_path),
        )
        assert first["success"] is True, first
        second = _post_office(
            client,
            "powerpoint",
            "write Second in title on slide 2",
            file_path=str(indexed_path),
        )
        assert second["success"] is True, second
        prs = Presentation(str(indexed_path))
        first_text = " ".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
        second_text = " ".join(shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text"))
        assert "first" in first_text.lower(), first_text
        assert "second" in second_text.lower(), second_text
        bullet_second = _post_office(
            client,
            "powerpoint",
            "on slide 2 add bullet Revenue increased",
            file_path=str(indexed_path),
        )
        assert bullet_second["success"] is True, bullet_second
        prs = Presentation(str(indexed_path))
        second_text = " ".join(shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text"))
        assert "revenue increased" in second_text.lower(), second_text

        duplicated = _post_office(
            client,
            "powerpoint",
            "duplicate slide 1",
            file_path=str(indexed_path),
        )
        assert duplicated["success"] is True, duplicated
        prs = Presentation(str(indexed_path))
        assert len(prs.slides) == 3
        duplicate_text = " ".join(shape.text for shape in prs.slides[2].shapes if hasattr(shape, "text"))
        assert "first" in duplicate_text.lower(), duplicate_text

        deleted = _post_office(
            client,
            "powerpoint",
            "delete slide 2",
            file_path=str(indexed_path),
        )
        assert deleted["success"] is True, deleted
        prs = Presentation(str(indexed_path))
        assert len(prs.slides) == 2
        remaining_text = " ".join(shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text"))
        assert "first" in remaining_text.lower(), remaining_text

        command_map.save_actions("excel", "create workbook", [{"action": "create_workbook"}])
        key, actions, score = command_map.get_cached_actions("excel", "create workbook with table")
        assert key is None and actions is None and score < 100

        handler = OpenAIHandler(api_key="")
        result = handler.interpret_result("excel", "add a complex chart")
        assert result.success is False
        assert result.error_code == "OPENAI_API_KEY_MISSING"
        parsed, _ = handler._parse_json('{"action":"create_workbook"}')
        assert normalize_actions(parsed) == [{"action": "create_workbook"}]
        try:
            handler._parse_json("{not valid json")
            raise AssertionError("Malformed JSON should fail")
        except OfficeActionError as exc:
            assert exc.error_code == "OPENAI_INVALID_JSON"

        unknown = _payload(client.post("/execute", json={"command": "open someunknownapp"}))
        assert unknown["success"] is False, unknown
        assert unknown["intent"] == "app_launch", unknown
        assert unknown.get("requires_manual_selection") is True, unknown

        invalid = _payload(client.post("/execute", json={"command": "dance please"}))
        assert invalid["success"] is False, invalid
        assert invalid["error_code"] == "UNKNOWN_COMMAND", invalid

        page = client.get("/")
        assert page.status_code == 200
        assert b"/static/reliability.js" in page.data

        frontend = client.get("/static/reliability.js")
        assert frontend.status_code == 200
        assert b"finally" in frontend.data
        assert b"Unexpected backend response" in frontend.data

    finally:
        if original_command_map is None:
            if COMMAND_MAP.exists():
                COMMAND_MAP.unlink()
        else:
            COMMAND_MAP.write_bytes(original_command_map)

    print("Office route smoke tests passed.")


if __name__ == "__main__":
    main()
