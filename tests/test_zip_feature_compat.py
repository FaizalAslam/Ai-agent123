import sys
import unittest
from pathlib import Path

VENV_SITE = Path(__file__).resolve().parent.parent / ".venv" / "Lib" / "site-packages"
if VENV_SITE.exists():
    sys.path.insert(0, str(VENV_SITE))
    try:
        import pptx  # noqa: F401
    finally:
        sys.path.remove(str(VENV_SITE))

from ai.openai_handler import OpenAIHandler
from executor.ppt_executor import PowerPointExecutor
from pptx import Presentation
from utils.office_actions import OfficeActionError, normalize_actions, validate_actions


class ZipFeatureCompatibilityTests(unittest.TestCase):
    def test_normalize_actions_accepts_legacy_parameter_shapes(self):
        self.assertEqual(
            normalize_actions({"action": "write_cell", "parameters": {"cell": "A1", "value": "Hello"}}),
            [{"action": "write_cell", "cell": "A1", "value": "Hello"}],
        )
        self.assertEqual(
            normalize_actions({"commands": [{"action": "write_cell", "params": {"cell": "B2", "value": 42}}]}),
            [{"action": "write_cell", "cell": "B2", "value": 42}],
        )
        self.assertEqual(
            normalize_actions([{"action": 'write_cell {"cell":"C3","value":"Done"}'}]),
            [{"action": "write_cell", "cell": "C3", "value": "Done"}],
        )

    def test_normalize_actions_rejects_bad_legacy_parameter_shapes(self):
        with self.assertRaises(OfficeActionError):
            normalize_actions({"action": "write_cell", "parameters": "cell=A1"})

    def test_openai_parser_accepts_legacy_command_containers(self):
        handler = OpenAIHandler(api_key="")
        actions, warnings, output_filename, ai_context = handler._parse_json(
            '{"commands":[{"action":"create_workbook"}],"output_filename":"Book.xlsx","context":{"table_range":"A1:C2"}}'
        )
        self.assertEqual(actions, [{"action": "create_workbook"}])
        self.assertEqual(warnings, [])
        self.assertEqual(output_filename, "Book.xlsx")
        self.assertEqual(ai_context, {"table_range": "A1:C2"})

        steps, _warnings, _filename, _context = handler._parse_json(
            '{"steps":[{"action":"create_document"}]}'
        )
        self.assertEqual(steps, [{"action": "create_document"}])

    def test_powerpoint_legacy_alias_actions_validate_and_execute(self):
        actions = validate_actions(
            "powerpoint",
            normalize_actions(
                [
                    {"action": "create_presentation"},
                    {"action": "add_slide", "title": "Original"},
                    {"action": "append_to_body", "slide_index": 1, "text": "Legacy body text"},
                    {"action": "add_speaker_notes", "slide_index": 1, "text": "Legacy notes"},
                    {"action": "find_replace_text", "find_text": "Legacy", "replace_text": "Updated"},
                    {"action": "apply_theme", "theme": "light_clean", "slide_index": 1},
                ]
            ),
        )
        executor = PowerPointExecutor(Presentation())
        results = [executor.run(action) for action in actions]
        self.assertTrue(all(result["status"] == "success" for result in results), results)
        slide_text = " ".join(shape.text for shape in executor.prs.slides[0].shapes if hasattr(shape, "text"))
        self.assertIn("Updated body text", slide_text)
        self.assertIn("Legacy notes", executor.prs.slides[0].notes_slide.notes_text_frame.text)


if __name__ == "__main__":
    unittest.main()
