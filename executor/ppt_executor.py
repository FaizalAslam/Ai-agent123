# modules/powerpoint_executor.py
import logging
import os

logger = logging.getLogger("OfficeAgent")


def _pt_color(hex_color):
    from pptx.dml.color import RGBColor
    named = {
        "red": "FF0000",
        "green": "00B050",
        "blue": "0070C0",
        "yellow": "FFFF00",
        "orange": "FFA500",
        "purple": "7030A0",
        "pink": "FF69B4",
        "black": "000000",
        "white": "FFFFFF",
        "gray": "808080",
        "grey": "808080",
        "navy": "000080",
        "darkblue": "00008B",
        "dark blue": "00008B",
    }
    raw = str(hex_color or "000000").strip().lower()
    h = named.get(raw, raw).lstrip("#")
    if len(h) == 8:
        h = h[-6:]
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class PowerPointExecutor:
    def __init__(self, prs):
        self.prs = prs

    def run(self, action_dict):
        action  = action_dict.get("action", "unknown")
        handler = getattr(self, f"_do_{action}", None)
        if not handler:
            logger.warning(f"PPT: Unknown action '{action}'")
            return {"status": "failed", "action": action, "message": f"Unknown action: {action}", "error_code": "UNKNOWN_ACTION"}
        try:
            handler(action_dict)
            return {"status": "success", "action": action, "message": ""}
        except Exception as e:
            logger.error(f"PPT action '{action}' failed: {e}")
            return {"status": "failed", "action": action, "message": str(e), "error_code": "ACTION_EXECUTION_ERROR"}


    def _do_create_presentation(self, p):
        from pptx import Presentation
        self.prs = Presentation()

    def _do_open_presentation(self, p):
        # Open/load is handled by server-level lifecycle.
        return

    # â”€â”€ Helpers â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _internal_slide_index(self, value=None, default=1):
        # External/user slide numbers are 1-based; python-pptx indexes are 0-based.
        try:
            external_index = int(value if value is not None else default)
        except (TypeError, ValueError):
            external_index = default
        return max(0, external_index - 1)

    def _slide(self, p):
        idx = self._internal_slide_index(p.get("slide_index"), default=1)
        if not self.prs.slides:
            self._do_add_slide({})
        return self.prs.slides[min(idx, len(self.prs.slides) - 1)]

    def _get_shape(self, slide, target):
        if not target:
            return slide.shapes[0] if slide.shapes else None
        for shape in slide.shapes:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph = str(shape.placeholder_format.type).upper()
                if target == "title"    and "TITLE"    in ph: return shape
                if target == "subtitle" and ("SUBTITLE" in ph or "BODY" in ph): return shape
                if target == "body"     and "BODY"     in ph: return shape
        return slide.shapes[0] if slide.shapes else None

    def _set_run_text(self, shape, text):
        if not shape or not shape.has_text_frame:
            return
        tf   = shape.text_frame
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
        else:
            para.add_run().text = text

    # â”€â”€ Slide Management â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_add_slide(self, p):
        layout_map = {
            "blank":         6,
            "title_only":    5,
            "title_content": 1,
            "two_content":   3,
        }
        idx    = layout_map.get(p.get("layout", "title_content"), 1)
        idx    = min(idx, len(self.prs.slide_layouts) - 1)
        layout = self.prs.slide_layouts[idx]
        self.prs.slides.add_slide(layout)

    def _do_delete_slide(self, p):
        idx = self._internal_slide_index(p.get("slide_index"), default=1)
        if idx < len(self.prs.slides):
            xml_slides = self.prs.slides._sldIdLst
            xml_slides.remove(xml_slides[idx])

    def _do_duplicate_slide(self, p):
        import copy
        template = self._slide(p)
        layout_idx = min(6, len(self.prs.slide_layouts) - 1)
        duplicate = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        # External/user slide index is 1-based; duplicate is appended to the
        # end while preserving the selected slide's visible shapes.
        for shape in list(duplicate.shapes):
            element = shape.element
            element.getparent().remove(element)
        for shape in template.shapes:
            duplicate.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")

    def _do_reorder_slide(self, p):
        from_idx = self._internal_slide_index(p.get("from_index"), default=1)
        to_idx   = self._internal_slide_index(p.get("to_index"), default=2)
        xml_list = self.prs.slides._sldIdLst
        item     = xml_list[from_idx]
        xml_list.remove(item)
        xml_list.insert(to_idx, item)

    def _do_hide_slide(self, p):
        self._slide(p)._element.set("show", "0")

    def _do_show_slide(self, p):
        self._slide(p)._element.set("show", "1")

    def _do_go_to_slide(self, p):
        logger.info(f"Navigate to slide {p.get('slide_index')} (requires win32com)")

    def _do_change_layout(self, p):
        layout_map = {
            "blank":         6,
            "title_only":    5,
            "title_content": 1,
            "two_content":   3,
        }
        idx    = layout_map.get(p.get("layout", "title_content"), 1)
        idx    = min(idx, len(self.prs.slide_layouts) - 1)
        slide  = self._slide(p)
        slide.slide_layout = self.prs.slide_layouts[idx]

    # â”€â”€ Text Content â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_set_slide_text(self, p):
        slide  = self._slide(p)
        target = p.get("target", "title")
        shape  = self._get_shape(slide, target)
        self._set_run_text(shape, p.get("text", ""))

    def _do_clear_slide_text(self, p):
        slide  = self._slide(p)
        target = p.get("target", "body")
        shape  = self._get_shape(slide, target)
        if shape and shape.has_text_frame:
            shape.text_frame.clear()

    def _do_add_bullet_point(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, "body")
        if shape and shape.has_text_frame:
            para       = shape.text_frame.add_paragraph()
            para.text  = p.get("text", "")
            para.level = 0

    def _do_add_numbered_point(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, "body")
        if shape and shape.has_text_frame:
            para       = shape.text_frame.add_paragraph()
            para.text  = p.get("text", "")
            para.level = 0

    def _do_set_speaker_notes(self, p):
        slide = self._slide(p)
        slide.notes_slide.notes_text_frame.text = p.get("text", "")

    # â”€â”€ Font & Style â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_set_font_size(self, p):
        from pptx.util import Pt
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(int(p["size"]))

    def _do_set_font_name(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = p["name"]

    def _do_set_font_color(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = _pt_color(p["color"])

    def _do_set_bold(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = p.get("bold", True)

    def _do_set_italic(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.italic = p.get("italic", True)

    def _do_set_underline(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.underline = p.get("underline", True)

    def _do_set_strikethrough(self, p):
        from pptx.oxml.ns import qn
        from lxml import etree
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr    = run._r.get_or_add_rPr()
                    strike = etree.SubElement(rPr, qn("a:strike"))
                    strike.set("val", "sngStrike")

    def _do_set_text_alignment(self, p):
        from pptx.enum.text import PP_ALIGN
        align_map = {
            "center":  PP_ALIGN.CENTER,
            "right":   PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
            "left":    PP_ALIGN.LEFT,
        }
        align = align_map.get(p.get("alignment", "left"), PP_ALIGN.LEFT)
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.alignment = align

    def _do_set_line_spacing(self, p):
        slide = self._slide(p)
        shape = self._get_shape(slide, p.get("target"))
        if shape and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para.line_spacing = float(p.get("spacing", 1.15))

    # â”€â”€ Background & Theme â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_set_bg_color(self, p):
        slide = self._slide(p)
        fill  = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _pt_color(p["color"])

    def _do_set_bg_image(self, p):
        from pptx.util import Inches
        path  = p.get("path", "")
        if not path or not os.path.exists(path):
            logger.warning(f"Background image not found: {path}")
            return
        slide = self._slide(p)
        slide.shapes.add_picture(
            path, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )

    def _do_set_bg_gradient(self, p):
        slide = self._slide(p)
        fill  = slide.background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = _pt_color(p.get("color1", "FFFFFF"))
        fill.gradient_stops[1].color.rgb = _pt_color(p.get("color2", "000000"))

    def _do_set_theme(self, p):
        logger.info(f"Theme '{p.get('theme')}' (requires win32com)")

    def _do_set_color_scheme(self, p):
        logger.info(f"Color scheme '{p.get('scheme')}' (requires win32com)")

    def _do_edit_slide_master(self, p):
        logger.info("Slide master edit (requires win32com)")

    def _do_add_logo(self, p):
        from pptx.util import Inches
        path = p.get("path", "")
        if not path or not os.path.exists(path):
            logger.warning(f"Logo not found: {path}")
            return
        for slide in self.prs.slides:
            slide.shapes.add_picture(path, Inches(0.1), Inches(0.1), Inches(1), Inches(0.5))

    # â”€â”€ Insert Objects â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_insert_image(self, p):
        from pptx.util import Inches
        path = p.get("path", "")
        if not path or not os.path.exists(path):
            logger.warning(f"Image not found: {path}")
            return
        slide = self._slide(p)
        slide.shapes.add_picture(path, Inches(1), Inches(1), Inches(4), Inches(3))

    def _do_resize_image(self, p):
        logger.info(f"Resize image to {p.get('width')}x{p.get('height')} (use shape index for targeting)")

    def _do_insert_shape(self, p):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        shape_map = {
            "RECTANGLE":    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            "OVAL":         MSO_AUTO_SHAPE_TYPE.OVAL,
            "RIGHT_ARROW":  MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            "TRIANGLE":     MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            "STAR_5_POINT": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            "PENTAGON":     MSO_AUTO_SHAPE_TYPE.PENTAGON,
            "DIAMOND":      MSO_AUTO_SHAPE_TYPE.DIAMOND,
        }
        shape_type = shape_map.get(p.get("shape_type", "RECTANGLE"), MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        slide      = self._slide(p)
        slide.shapes.add_shape(shape_type, Inches(1), Inches(1), Inches(2), Inches(1.5))

    def _do_insert_text_box(self, p):
        from pptx.util import Inches
        slide = self._slide(p)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = p.get("text", "")

    def _do_insert_table(self, p):
        from pptx.util import Inches
        slide = self._slide(p)
        slide.shapes.add_table(
            int(p.get("rows", 3)), int(p.get("cols", 3)),
            Inches(1), Inches(1), Inches(6), Inches(3)
        )

    def _do_insert_chart(self, p):
        from pptx.util import Inches
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        chart_map = {
            "bar":     XL_CHART_TYPE.BAR_CLUSTERED,
            "column":  XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line":    XL_CHART_TYPE.LINE,
            "pie":     XL_CHART_TYPE.PIE,
            "area":    XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }
        chart_type = chart_map.get(p.get("chart_type", "bar"), XL_CHART_TYPE.BAR_CLUSTERED)
        chart_data = ChartData()
        chart_data.categories = ["A", "B", "C"]
        chart_data.add_series("Series 1", (1, 2, 3))
        slide = self._slide(p)
        slide.shapes.add_chart(chart_type, Inches(1), Inches(1), Inches(6), Inches(4), chart_data)

    def _do_insert_video(self, p):
        logger.info(f"Video '{p.get('path')}' (requires win32com)")

    def _do_insert_audio(self, p):
        logger.info(f"Audio '{p.get('path')}' (requires win32com)")

    def _do_insert_icon(self, p):
        logger.info(f"Icon '{p.get('icon_name')}' (requires win32com / Microsoft 365)")

    def _do_insert_smartart(self, p):
        logger.info(f"SmartArt '{p.get('type')}' (requires win32com)")

    def _do_insert_hyperlink(self, p):
        from pptx.util import Inches
        slide = self._slide(p)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
        run   = txBox.text_frame.paragraphs[0].add_run()
        run.text             = p.get("text", p.get("url", "Link"))
        run.hyperlink.address = p.get("url", "")

    # â”€â”€ Object Positioning â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_align_object_left(self, p):
        logger.info("Align object left (requires win32com)")

    def _do_align_object_center(self, p):
        logger.info("Align object center (requires win32com)")

    def _do_align_object_right(self, p):
        logger.info("Align object right (requires win32com)")

    def _do_align_object_top(self, p):
        logger.info("Align object top (requires win32com)")

    def _do_align_object_bottom(self, p):
        logger.info("Align object bottom (requires win32com)")

    def _do_align_object_middle(self, p):
        logger.info("Align object middle (requires win32com)")

    def _do_align_object(self, p):
        logger.info(f"Align object '{p.get('alignment')}' (requires win32com)")

    def _do_group_objects(self, p):
        logger.info("Group objects (requires win32com)")

    def _do_ungroup_objects(self, p):
        logger.info("Ungroup objects (requires win32com)")

    def _do_bring_forward(self, p):
        logger.info("Bring forward (requires win32com)")

    def _do_send_backward(self, p):
        logger.info("Send backward (requires win32com)")

    def _do_rotate_object(self, p):
        slide = self._slide(p)
        for shape in slide.shapes:
            shape.rotation = float(p.get("angle", 90))

    def _do_flip_object(self, p):
        logger.info(f"Flip '{p.get('direction')}' (requires win32com)")

    # â”€â”€ Transitions & Animations â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_set_transition(self, p):
        logger.info(f"Transition '{p.get('transition')}' on slide {p.get('slide_index')} (requires win32com)")

    def _do_set_transition_speed(self, p):
        logger.info(f"Transition speed '{p.get('speed')}' (requires win32com)")

    def _do_remove_transition(self, p):
        from pptx.oxml.ns import qn
        slide  = self._slide(p)
        timing = slide._element.find(qn("p:timing"))
        if timing is not None:
            slide._element.remove(timing)

    def _do_apply_transition_all(self, p):
        logger.info(f"Apply transition '{p.get('transition')}' to all slides (requires win32com)")

    def _do_set_animation(self, p):
        logger.info(f"Animation '{p.get('animation')}' on slide {p.get('slide_index')} (requires win32com)")

    def _do_set_animation_delay(self, p):
        logger.info(f"Animation delay {p.get('delay')}s (requires win32com)")

    def _do_remove_animation(self, p):
        logger.info("Remove animation (requires win32com)")

    def _do_set_auto_advance(self, p):
        from pptx.util import Pt
        slide = self._slide(p)
        slide.slide_layout.slide_master.slide_layouts
        logger.info(f"Auto-advance after {p.get('seconds')}s (requires win32com)")

    # â”€â”€ Slide Settings â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_set_slide_size(self, p):
        from pptx.util import Inches
        self.prs.slide_width  = Inches(float(p.get("width",  13.33)))
        self.prs.slide_height = Inches(float(p.get("height",  7.5)))

    def _do_set_header_footer(self, p):
        logger.info(f"Header '{p.get('header')}' / Footer '{p.get('footer')}' (requires win32com)")

    def _do_add_slide_number(self, p):
        logger.info("Slide number (requires win32com)")

    # â”€â”€ Presentation Tools â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_start_slideshow(self, p):
        logger.info("Start slideshow (requires win32com)")

    def _do_end_slideshow(self, p):
        logger.info("End slideshow (requires win32com)")

    def _do_print_handouts(self, p):
        logger.info("Print handouts (requires win32com)")

    def _do_spell_check(self, p):
        logger.info("Spell check (requires win32com)")

    # â”€â”€ Save / Export â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _do_save_presentation(self, p):
        path = getattr(self.prs, "_path", "output.pptx")
        self.prs.save(path)

    def _do_save_presentation_as(self, p):
        filename = p.get("filename", "output")
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        self.prs.save(filename)

    def _do_close_presentation(self, p):
        logger.info("Presentation closed (session ended)")

    def _do_export_pdf(self, p):
        logger.info(f"PDF export to '{p.get('path', 'output.pdf')}' (requires win32com)")

    def _do_undo(self, p):
        logger.info("Undo (requires win32com)")

    def _do_redo(self, p):
        logger.info("Redo (requires win32com)")
